
from flask import Flask, request, jsonify
import google.generativeai as genai
from werkzeug.utils import secure_filename
import os
from PyPDF2 import PdfReader
import docx
import pptx
import subprocess

app = Flask(__name__)

# Set up Google Gemini API
API_KEY = "AIzaSyA8cpAppBZIkbubvymxJwv52-1i3lVbnnE"  # Replace with actual API Key
genai.configure(api_key=API_KEY)

# Function to extract text from PDF
def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

# Function to extract text from DOCX
def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

# Function to extract text from PPTX
def extract_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to extract text from TXT
def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()
    return text

# Function to extract text from DOC (legacy format)
def extract_text_from_doc(file_path):
    try:
        # Convert `.doc` to `.docx` using unoconv or LibreOffice
        docx_file_path = file_path.replace('.doc', '.docx')
        subprocess.run(['unoconv', '-f', 'docx', file_path], check=True)
        
        # Extract text using python-docx
        return extract_text_from_docx(docx_file_path)
    except Exception as e:
        return f"Failed to process DOC file: {str(e)}"

# Function to generate a professional summary using AI
# Function to generate a professional summary using AI
def generate_summary(text):
    model = genai.GenerativeModel("gemini-1.5-flash")
    prompt = f"""
    Please summarize the following text and present it in a proposal format with clear sections. 
    **Important:** Please provide a 4 to 5-page proposal with detailed descriptions. Kindly avoid using any tabular format.
    Please write a proposal that sounds natural and uses simple, everyday language, but still includes professional terms where needed to make the proposal more effective. I want the descriptions to be unique and impactful. Kindly write the proposal in a way that is easy to understand, while maintaining a professional tone.

    For each section, provide comprehensive explanations:

    1. **Introduction**: Provide a brief introduction to the project, including the context and goals of the project. Explain why the project is important and what problem it aims to solve.
    
    2. **Project Scope**: Outline the main features and objectives of the project. Describe what is in scope and what is out of scope. Include the intended audience and the primary deliverables.

    3. **Project Timeline**: Define the project phases, including specific working days and buffer time. Provide a detailed breakdown of the timeline, including any important milestones and deadlines.

    4. **Technology Stack**: Specify the technologies to be used in the project. Provide details about each technology, including why it was chosen and how it fits into the overall architecture of the project.

    5. **Cost**: Provide an estimated cost breakdown for the project. Include a detailed explanation of the costs involved, such as development, testing, deployment, maintenance, and any other associated costs.

    6. **Our Workflow**: Outline the development and testing process in detail. Describe how the project will be managed, what methodology will be used (e.g., Agile, Scrum), and the steps involved in each phase of development.

    7. **Project Deliverables**: List the major deliverables after the project's completion, including software, documentation, and any other assets. Describe the quality standards and acceptance criteria for each deliverable.

    8. **Development Phase**: Describe the entire development phase in detail, from initial development to final code review and testing. Include information on how features will be implemented and tested in each sprint or phase.

    9. **Deployment**: Provide a detailed explanation of the deployment process, including how the project will be deployed to production. Include steps for testing, staging, and the final deployment. Also, describe any post-deployment monitoring or support that will be required.

    10. **Conclusion**: Summarize the entire project, reiterate the key points, and provide a clear overview of the next steps, including any follow-up actions or meetings needed to kick off the project.
     Sincerely,
        [Your name ]
        [Your Post]
        [Company name]
        [Email-ID]
        [Mobile No]
    
    Please ensure that the text follows a professional and formal tone suitable for a business proposal.
    
    **Text:**
    {text}
    """
    response = model.generate_content([prompt])
    return response.text


# Route to handle file uploads and text input
@app.route('/summarize', methods=['POST'])
def summ():
    file = request.files.get('file')
    text_input = request.form.get('text')

    # Ensure the 'uploads' directory exists
    upload_folder = 'uploads'
    if not os.path.exists(upload_folder):
        os.makedirs(upload_folder)

    if file:
        filename = secure_filename(file.filename)
        file_path = os.path.join(upload_folder, filename)
        try:
            file.save(file_path)
        except Exception as e:
            return jsonify({"message": f"Failed to save file: {str(e)}"}), 500
        
        file_ext = filename.split('.')[-1].lower()
        
        if file_ext == 'pdf':
            text = extract_text_from_pdf(file_path)
        elif file_ext == 'docx':
            text = extract_text_from_docx(file_path)
        elif file_ext == 'pptx':
            text = extract_text_from_pptx(file_path)
        elif file_ext == 'txt':
            text = extract_text_from_txt(file_path)
        elif file_ext == 'doc':
            text = extract_text_from_doc(file_path)
        else:
            return jsonify({"message": "Unsupported file format"}), 400
    elif text_input:
        text = text_input
    else:
        return jsonify({"message": "No file or text input provided"}), 400

    summary = generate_summary(text)
    return jsonify({"summary": summary})

if __name__ == '__main__':
    app.run(debug=True)

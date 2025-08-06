from flask import Flask, request, jsonify , Blueprint
import json
import google.generativeai as genai
import io
import docx
import pptx
from PyPDF2 import PdfReader
import re
from config import API_KEYS  # Import the API keys
import logging
from collections import Counter
import pandas as pd
import numpy as np
import pickle
from fuzzywuzzy import process  # For better module name matching
import random
import os
import chardet
from sklearn.linear_model import LinearRegression

app = Flask(__name__)

# Configure logging
logging.basicConfig(level=logging.INFO)

# Use API keys from config
current_key_index = 0

def reset_api_keys():
    """Reset all API keys to active status (status = True)."""
    for key_info in API_KEYS:
        key_info["status"] = True
    logging.info("All API keys have been reset to active.")

def configure_api_key():
    """Configure the first active API key."""
    global current_key_index
    active_keys = [key for key in API_KEYS if key["status"]]
    
    if not active_keys:
        reset_api_keys()  # Reset all keys if none are active
        active_keys = API_KEYS
    
    # Set the first active key
    for i, key_info in enumerate(API_KEYS):
        if key_info["status"]:
            current_key_index = i
            genai.configure(api_key=key_info["key"])
            logging.info(f"Configured API Key: {key_info['key']}")
            return

    # If no active key is found after reset (unlikely scenario)
    logging.error("No active API keys available.")
    raise ValueError("No active API keys available.")

def rotate_api_key():
    """Rotate to the next available API key."""
    global current_key_index
    API_KEYS[current_key_index]["status"] = False  # Mark the current key as inactive
    current_key_index = (current_key_index + 1) % len(API_KEYS)

    active_keys = [key for key in API_KEYS if key["status"]]

    if not active_keys:
        reset_api_keys()  # Reset all keys if none are active
        active_keys = API_KEYS
    
    # Set the next available active key
    for i, key_info in enumerate(API_KEYS):
        if key_info["status"]:
            current_key_index = i
            genai.configure(api_key=key_info["key"])
            logging.info(f"Rotated to API Key: {key_info['key']}")
            return

    # If no active key is found after reset (unlikely scenario)
    logging.error("No active API keys available after rotation.")
    raise ValueError("No active API keys available.")

# Initial configuration - commented out to avoid hanging during import
# try:
#     configure_api_key()
# except ValueError as e:
#     logging.error(f"API Key Configuration Error: {e}")

def extract_text_from_pdf(file):
    """Extract text from a PDF file."""
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_docx(file):
    """Extract text from a DOCX file."""
    doc = docx.Document(file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_pptx(file):
    """Extract text from a PPTX file."""
    prs = pptx.Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def detect_languages(text):
    """Detect frontend and backend languages from the given text."""
    try:
        # Ensure text is properly encoded
        if isinstance(text, bytes):
            text = text.decode('utf-8', errors='ignore')
        elif not isinstance(text, str):
            text = str(text)
        
        # Configure API key first
        configure_api_key()
        model = genai.GenerativeModel("gemini-1.5-flash")
        prompt = f"""
        Analyze the project description below and extract all mentioned programming languages, frameworks, and databases used for frontend and backend development. 

        **Important Instructions:**
        - If specific technologies are mentioned, use those exact names
        - If no technologies are mentioned, suggest the most suitable options based on the project type
        - For e-commerce projects, suggest: Frontend: React, Angular, Vue.js | Backend: Node.js, PHP, Python
        - For CRM projects, suggest: Frontend: React, Angular | Backend: Node.js, Python, Java
        - For general web projects, suggest: Frontend: React, HTML, CSS | Backend: Node.js, PHP
        - Always provide at least one frontend and one backend technology

        **Format your response exactly as follows:**
        Frontend: [technology1, technology2]
        Backend: [technology1, technology2]

        **Project Description:**
        {text}
        """

        response = model.generate_content([prompt])
        response_text = response.text.strip()
        # Ensure response text is properly encoded
        if isinstance(response_text, bytes):
            response_text = response_text.decode('utf-8', errors='ignore')
        logging.info("Response Text: %s", response_text)

        if not response_text:
            logging.error("Received empty response from Gemini API.")
            return [], []

        frontend_match = re.search(r'Frontend:\s*(.*)', response_text, re.IGNORECASE)
        backend_match = re.search(r'Backend:\s*(.*)', response_text, re.IGNORECASE)

        detected_frontend = []
        detected_backend = []

        if frontend_match:
            detected_frontend = [tech.strip() for tech in frontend_match.group(1).split(',')]
        else:
            logging.warning("No frontend technologies detected.")

        if backend_match:
            detected_backend = [tech.strip() for tech in backend_match.group(1).split(',')]
        else:
            logging.warning("No backend technologies detected.")

        logging.info(f"Detected Frontend Technologies: {', '.join(detected_frontend) if detected_frontend else 'None detected'}")
        logging.info(f"Detected Backend Technologies: {', '.join(detected_backend) if detected_backend else 'None detected'}")

        return detected_frontend, detected_backend

    except Exception as e:
        logging.error(f"Error detecting languages: {e}")
        return [], []
    

# Model loading function to avoid import issues
def load_models():
    """Load trained models and encoders"""
    model_path = "model/"
    dataset_path = "model/dataset.json"
    
    try:
        logging.info("Loading time prediction model...")
        time_model = pickle.load(open(model_path + "time_prediction_model.pkl", "rb"))
        logging.info("Loading cost prediction model...")
        cost_model = pickle.load(open(model_path + "cost_prediction_model.pkl", "rb"))
        logging.info("Loading module name encoder...")
        label_enc_module = pickle.load(open(model_path + "module_name_encoder.pkl", "rb"))
        logging.info("Loading language encoder...")
        label_enc_lang = pickle.load(open(model_path + "language_encoder.pkl", "rb"))
        
        # Load dataset from JSON
        try:
            logging.info("Loading dataset...")
            with open(dataset_path, "r") as f:
                dataset = json.load(f)
            logging.info(f"Dataset loaded with {len(dataset)} entries")
        except (FileNotFoundError, json.JSONDecodeError) as e:
            logging.warning(f"Dataset loading failed: {e}")
            dataset = []
            
        logging.info("All models loaded successfully")
        return time_model, cost_model, label_enc_module, label_enc_lang, dataset
    except Exception as e:
        logging.error(f"Error loading models: {e}")
        import traceback
        logging.error(f"Full traceback: {traceback.format_exc()}")
        raise e


def get_closest_module(module_name, user_id=None, label_enc_module=None):
    """Find the closest module match using fuzzy matching, with optional user-specific dataset."""
    # Use provided label encoder or load models if not provided
    if label_enc_module is None:
        time_model, cost_model, label_enc_module, label_enc_lang, dataset = load_models()
    
    if user_id:
        user_directory = os.path.join('user_data', str(user_id))
        user_file_path = os.path.join(user_directory, f'{user_id}.json')
        try:
            with open(user_file_path, "r") as f:
                user_dataset = json.load(f)
            existing_modules = [module["Module Name"].lower() for module in user_dataset]
        except (FileNotFoundError, json.JSONDecodeError):
            logging.warning(f"User-specific dataset not found for user {user_id}, using default dataset.")
            existing_modules = label_enc_module.classes_
    else:
        existing_modules = label_enc_module.classes_

    # Clean the module name for better matching
    cleaned_module = module_name.lower().replace("-", "").strip()
    
    # Try exact match first
    if cleaned_module in existing_modules:
        return cleaned_module
    
    # Try partial matches
    for existing_module in existing_modules:
        if cleaned_module in existing_module or existing_module in cleaned_module:
            logging.info(f"Partial match found: '{cleaned_module}' -> '{existing_module}'")
            return existing_module
    
    # Use fuzzy matching with lower threshold
    best_match, confidence = process.extractOne(cleaned_module, existing_modules) if len(existing_modules) > 0 else (None, 0)
    
    # Lower the confidence threshold to get more matches
    if confidence >= 60:  # Reduced from 80 to 60
        logging.info(f"Fuzzy match found: '{cleaned_module}' -> '{best_match}' (confidence: {confidence})")
        return best_match
    
    # If still no match, try to map common patterns
    common_mappings = {
        'otp verification': 'user authentication',
        'add car': 'product management',
        'edit car': 'product management', 
        'delete car': 'product management',
        'car information management': 'inventory management',
        'car search': 'search',
        'car filters': 'search',
        'car details': 'product catalog',
        'car listing': 'product catalog',
        'car management': 'inventory management',
        'insurance policy': 'content management',
        'insurance management': 'content management',
        'loan management': 'order management',
        'payment gateway': 'payment gateway',
        'transaction history': 'payment history',
        'push notifications': 'notification system',
        'data synchronization': 'data analytics',
        'admin dashboard': 'dashboard',
        'admin user management': 'user management',
        'admin car management': 'product management',
        'admin insurance management': 'content management',
        'admin loan management': 'order management',
        'report generation': 'report generation',
        'analytics dashboard': 'analytics dashboard',
        'user feedback system': 'feedback',
        'customer support': 'help',
        'security system': 'security management',
        'data backup': 'file management',
        'version control': 'api integrations',
        'api integration': 'api integrations',
        'database management': 'content management',
        'server management': 'user management',
        'maintenance and updates': 'settings',
        'user help section': 'help',
        'faqs': 'faq',
        'privacy policy': 'about',
        'terms of service': 'about',
        'instagram integration': 'social media integration',
        'whatsapp integration': 'social media integration',
        'location services': 'cms integrations',
        'my cars screen': 'dashboard',
        'home screen': 'dashboard',
        'cars section': 'product catalog',
        'insurance section': 'content management',
        'loans section': 'order management',
        'bottom navigation bar': 'navigation'
    }
    
    # Try common mappings
    for pattern, mapped_module in common_mappings.items():
        if pattern in cleaned_module:
            if mapped_module in existing_modules:
                logging.info(f"Pattern mapping: '{cleaned_module}' -> '{mapped_module}'")
                return mapped_module
    
    # Last resort: return the original module name
    logging.warning(f"No match found for '{cleaned_module}', using original name")
    return module_name

def predict_time_cost(module_name, language):
    """Predict time and cost using the model from new_test.py, but with automatic inputs."""
    
    # Load models
    time_model, cost_model, label_enc_module, label_enc_lang, dataset = load_models()
    
    # Convert to lowercase for consistency
    module_name = module_name.lower()
    language = language.lower()

    # Find closest module match
    best_match = get_closest_module(module_name)

    # Encode module name
    module_encoded = label_enc_module.transform([best_match])[0] if best_match in label_enc_module.classes_ else len(label_enc_module.classes_)

    # Encode technology
    lang_encoded = label_enc_lang.transform([language])[0] if language in label_enc_lang.classes_ else len(label_enc_lang.classes_)

    # Predict time and cost
    predicted_time = time_model.predict([[module_encoded, lang_encoded]])[0]
    predicted_cost = cost_model.predict([[module_encoded, lang_encoded]])[0]

    # Determine Complexity
    complexity = "Small" if predicted_time <= 80 else "Medium" if predicted_time <= 20 else "Large"

    # Save prediction to JSON dataset
    new_entry = {
        "Module Name": module_name,
        "Technology": language,
        "Time": round(predicted_time),
        "Cost": round(predicted_cost),
        "Complexity": complexity
    }

    # Append new data to JSON and save
    # dataset.append(new_entry)
    # with open(dataset_path, "w") as f:
    #     json.dump(dataset, f, indent=4)

    return {
        "Module": module_name,
        "Technology": language,
        "Predicted Time (hrs)": round(predicted_time, 2),
        "Predicted Cost (INR)": round(predicted_cost, 2),
        "Complexity": complexity
    }

import re

def clean_module_names(module_text):
    """Extracts and cleans module names from generated text."""
    module_names = []
    lines = module_text.strip().split("\n")

    for line in lines:
        # Remove numbering (e.g., "1. login" → "login")
        cleaned_line = re.sub(r"^\d+\.\s*", "", line).strip()
        
        # Skip empty lines and section headers (e.g., "**core website functionalities:**")
        if cleaned_line and not cleaned_line.startswith("**"):
            module_names.append(cleaned_line)

    return module_names


def generate_module_details(text):
    """Generate module details based on the provided text."""
    max_retries = 3
    retry_count = 0
    
    # Ensure text is properly encoded
    if isinstance(text, bytes):
        text = text.decode('utf-8', errors='ignore')
    elif not isinstance(text, str):
        text = str(text)
    
    while retry_count < max_retries:
        try:
            # Configure API key first
            configure_api_key()
            detected_frontend, detected_backend = detect_languages(text)

            model = genai.GenerativeModel("gemini-1.5-flash")
            website_keywords = [
                "website", "user authentication", "content management", "admin panel",
                "product catalog", "shopping cart", "checkout", "payment gateway", 
                "user dashboard", "profile management", "order tracking", "search functionality", 
                "feedback system", "multi-language support", "SEO optimization"
            ]

            detected_features = [kw for kw in website_keywords if re.search(r'\b' + re.escape(kw) + r'\b', text, re.IGNORECASE)]
            admin_keywords = ["admin panel", "admin dashboard", "admin"]
            admin_side_detected = any(re.search(r'\b' + re.escape(kw) + r'\b', text, re.IGNORECASE) for kw in admin_keywords)

            prompt_modules = f"""
Analyze the provided project description and extract the following details with high accuracy:

1. **Client Name**: Identify the name of the client, company, or organization mentioned in the text.  
   - If explicitly stated, extract the exact name.  
   - If not mentioned directly, infer a plausible name based on the context (e.g., company branding, domain-specific names, or text hints).  
   - If no relevant name is found, return "Not Specified".  

2. **Project Title**: Extract the name or title of the project.  
   - If explicitly mentioned, use the exact wording.  
   - If the title is unclear, infer a meaningful and concise title based on the project's description.  
   - If no clear project title is found, return "Not Specified".  

3. **List of Module Names**: Extract all relevant module names mentioned or implied in the text.  
   - Return module names as a structured list (without descriptions).  
   - If module names are implied, infer logical names based on the project type.  
   - Ensure the list contains at least **40 module names** (if applicable) related to the project's domain.

### **Important Considerations:**
- The response should **only** contain the required details in the structured format below.  
- Avoid explanations, additional text, or generic assumptions.  
- Use **formal and structured module names** (e.g., "User Authentication", "Admin Dashboard", "Payment Gateway").  
- If the project is related to **eCommerce, CRM, ERP, or SaaS**, generate modules relevant to these domains.  
equired format.

            **Format your response as follows:**
            Client Name: [client name]
            Project Title: [project title]
            Modules:
            - [module name 1]
            - [module name 2]
            - [module name 3]
            ...
            
            **Project Description:**

            **Text:**
            {text}
            Please generate at least 40 module names. If the text relates to a specific domain like eCommerce or CRM, include relevant module names.
            """

            if detected_features:
                prompt_modules += f"""
                The following website features were detected in the text: {', '.join(detected_features)}.
                Please focus on generating modules related to website functionalities, such as 'user authentication', 'product catalog', 'shopping cart', 'payment gateway', 'admin panel', etc.
                """
            if admin_side_detected:
                prompt_modules += """
                The following admin-related features were detected. Please focus on generating modules relevant to the admin side, such as 'admin dashboard', 'user management', 'order management', 'product management', 'report generation', etc.
                """

            response_modules = model.generate_content([prompt_modules])

            # print(response_modules)


            # Parse the response
            response_text = response_modules.text.strip()
            # Ensure response text is properly encoded
            if isinstance(response_text, bytes):
                response_text = response_text.decode('utf-8', errors='ignore')
            lines = response_text.split("\n")

            client_name = "Not Specified"
            project_title = "Not Specified"
            module_names = []

            # Extract client name, project title, and modules from the response
            parsing_modules = False
            for line in lines:
                line = line.strip()
                if line.startswith("Client Name:"):
                    client_name = line.replace("Client Name:", "").strip()
                elif line.startswith("Project Title:"):
                    project_title = line.replace("Project Title:", "").strip()
                elif line.startswith("Modules:"):
                    parsing_modules = True
                elif parsing_modules and line.startswith("-"):
                    module_names.append(line.replace("-", "").strip())

            # Clean module names (optional, if needed beyond initial cleaning)
            cleaned_modules = clean_module_names("\n".join([f"- {m}" for m in module_names]))

            # print(cleaned_modules)

            

            module_details = []

            for i, module in enumerate(cleaned_modules):
                if i == 0 or i == len(cleaned_modules) - 1:
                    continue

                # Debug: Log what we have for technology detection
                logging.info(f"Module: {module}")
                logging.info(f"Detected Frontend: {detected_frontend}")
                logging.info(f"Detected Backend: {detected_backend}")
                
                # Fix: Use first backend technology or default to a common one
                backend_tech = detected_backend[0] if detected_backend else "PHP"
                logging.info(f"Using backend technology: {backend_tech}")
                
                time_cost_prediction = predict_time_cost(module, backend_tech)


                module_details.append({
                    "index": i,
                    "module_name": module,
                    "frontend": detected_frontend,
                    "backend": detected_backend,
                    "predicted_time": time_cost_prediction["Predicted Time (hrs)"],
                    "predicted_cost": time_cost_prediction["Predicted Cost (INR)"],
                    "complexity": time_cost_prediction["Complexity"] 
                })

            if module_details:
                return {
                    "client_name": client_name,
                    "project_title": project_title,
                    "modules": module_details
                }
            else:
                logging.warning("Module details are empty. Retrying...")
                retry_count += 1
                continue

        except Exception as e:
            if "429" in str(e):
                rotate_api_key()
                retry_count += 1
                continue
            else:
                logging.error(f"Error during module generation: {e}")
                retry_count += 1
                if retry_count >= max_retries:
                    raise e
                continue
    
    # If we get here, all retries failed
    raise Exception("Failed to generate module details after maximum retries")


#  api call mate no fuunction amathi badho response html file ma jashe
def pruposal_module():
    print("pruposal_module function called - this should not happen!")
    logging.info("pruposal_module function called - this should not happen!")
    user_id = request.form.get("user_id")
    file = request.files.get("file")
    text_input = request.form.get("text")

    if not file and not text_input:
        return jsonify({"error": "No input text or file provided"}), 400

    try:
        text = file.read().decode("utf-8") if file else text_input

        # Determine dataset and model based on user ID
        if not user_id:
            # Use default dataset and model
            dataset = load_default_dataset()
            model = load_default_model()
        else:
            user_data_path = f"user_data/{user_id}/"
            if not os.path.exists(user_data_path):
                os.makedirs(user_data_path)

            # Check if user-specific dataset and model exist
            user_dataset_path = os.path.join(user_data_path, "dataset.csv")
            user_model_path = os.path.join(user_data_path, "model.pkl")

            if not os.path.exists(user_dataset_path) or not os.path.exists(user_model_path):
                # First-time user: use default dataset and model
                dataset = load_default_dataset()
                model = load_default_model()
                # Save initial response and train new model
                response = generate_module_details(text)
                save_response(user_id, response)
                train_and_save_model(user_id, dataset, response)
            else:
                # Returning user: use user-specific dataset and model
                dataset = load_user_dataset(user_id)
                model = load_user_model(user_id)
                response = generate_module_details(text)
                update_user_data(user_id, response)
                train_and_save_model(user_id, dataset, response)

        return jsonify({"modules": response}), 200

    except Exception as e:
        if "429" in str(e):
            rotate_api_key()
            return jsonify({"error": "Rate limit exceeded. Retrying with new API key."}), 429
        else:
            logging.error(f"Error: {str(e)}")
            return jsonify({"error": "An unknown error occurred."}), 500

# Helper functions to load datasets and models, save responses, and train models

def load_default_dataset():
    # Load the default dataset
    return pd.read_csv("model/dataset.csv")

def load_default_model():
    # Load the default pre-trained model
    with open("model/default_model.pkl", "rb") as f:
        return pickle.load(f)

def load_user_dataset(user_id):
    # Load the user-specific dataset
    return pd.read_csv(f"user_data/{user_id}/dataset.csv")

def load_user_model(user_id):
    # Load the user-specific model
    with open(f"user_data/{user_id}/model.pkl", "rb") as f:
        return pickle.load(f)

def save_response(user_id, response):
    # Save the response to a user-specific JSON file
    user_directory = os.path.join('user_data', str(user_id))
    os.makedirs(user_directory, exist_ok=True)
    user_file_path = os.path.join(user_directory, f'{user_id}.json')
    with open(user_file_path, "w") as f:
        json.dump(response, f, indent=4)

def update_user_data(user_id, response):
    # Update the user-specific JSON file with new data
    user_directory = os.path.join('user_data', str(user_id))
    os.makedirs(user_directory, exist_ok=True)
    user_file_path = os.path.join(user_directory, f'{user_id}.json')
    if os.path.exists(user_file_path):
        with open(user_file_path, "r") as f:
            existing_data = json.load(f)
        existing_data.extend(response)
        with open(user_file_path, "w") as f:
            json.dump(existing_data, f, indent=4)
    else:
        save_response(user_id, response)

def train_and_save_model(user_id, dataset, response):
    # Train a new model using the dataset and response, then save it
    # This is a placeholder for the actual training logic
    new_model = train_model(dataset, response)
    with open(f"user_data/{user_id}/model.pkl", "wb") as f:
        pickle.dump(new_model, f)

def train_model(dataset, response):
    # Placeholder function for training a model
    # Implement the actual training logic here
    return LinearRegression().fit(dataset, response)

def module_details_api():
    """API endpoint to get module details."""
    logging.info("module_details_api function called - this is correct!")
    
    # Load default models first
    time_model, cost_model, label_enc_module, label_enc_lang, dataset = load_models()
    
    # Debug: Log all request data
    logging.info(f"Request method: {request.method}")
    logging.info(f"Request headers: {dict(request.headers)}")
    logging.info(f"Request form data: {dict(request.form)}")
    logging.info(f"Request files: {dict(request.files)}")
    logging.info(f"Request JSON: {request.get_json() if request.is_json else 'Not JSON'}")
    
    # Check if request is JSON or form data
    if request.is_json:
        # Handle JSON request format
        data = request.get_json()
        text_input = data.get("text")
        user_id = data.get("user_id")
        file = None  # No file support in JSON format
    else:
        # Handle form data format (backward compatibility)
        file = request.files.get("file")
        text_input = request.form.get("text")
        user_id = request.form.get("user_id")

    # Debug: Log what we received
    logging.info(f"Received request - file: {file}, text: {text_input[:100] if text_input else 'None'}, user_id: {user_id}")

    if not file and not text_input:
        logging.error("No input text or file provided")
        return jsonify({"error": "No input text or file provided"}), 400

    try:
        if file:
            raw_data = file.read()
            result = chardet.detect(raw_data)
            encoding = result['encoding'] if result['encoding'] else 'latin1'
            text = raw_data.decode(encoding)
        else:
            text = text_input

        # Generate module details
        response = generate_module_details(text)
        logging.info(f"Generated module details: {response}")

        # If user_id is provided, save the response
        if user_id:
            # Define user-specific directory and file path
            user_directory = os.path.join('user_data', str(user_id))
            os.makedirs(user_directory, exist_ok=True)
            user_file_path = os.path.join(user_directory, f'{user_id}.json')
            model_directory = os.path.join(user_directory, 'models')

            # Check if user has existing models
            if os.path.exists(model_directory):
                # Load user-specific models
                user_time_model = pickle.load(open(os.path.join(model_directory, "time_prediction_model.pkl"), "rb"))
                user_cost_model = pickle.load(open(os.path.join(model_directory, "cost_prediction_model.pkl"), "rb"))
                logging.info(f"Loaded user-specific models for user {user_id}.")
            else:
                # Use default models for new users
                time_model, cost_model, label_enc_module, label_enc_lang, dataset = load_models()
                user_time_model = time_model
                user_cost_model = cost_model
                logging.info(f"Using default models for user {user_id}.")

            # Load existing user dataset
            try:
                with open(user_file_path, "r") as f:
                    user_dataset = json.load(f)
            except (FileNotFoundError, json.JSONDecodeError):
                user_dataset = []

            logging.info(f"Existing user dataset before update: {user_dataset}")

            # Update or replace module details
            for module in response["modules"]:
                module_name = module.get("module_name")
                
                # Fix: Handle empty backend array properly
                backend_array = module.get("backend", [])
                if backend_array and len(backend_array) > 0:
                    technology = backend_array[0]
                    # Clean technology format - remove brackets if present
                    if isinstance(technology, str):
                        technology = technology.strip()
                        if technology.startswith('[') and technology.endswith(']'):
                            technology = technology[1:-1].strip()
                else:
                    technology = "PHP"  # Default to PHP if no backend detected
                
                logging.info(f"Processing module: {module_name}, Technology: {technology}")

                # Clean module name - remove dashes and extra spaces
                cleaned_module_name = module_name.replace("-", "").strip()
                
                # Use fuzzy matching to find the best module match
                best_match = get_closest_module(cleaned_module_name, user_id, label_enc_module)
                logging.info(f"Original: '{cleaned_module_name}' -> Best Match: '{best_match}'")
                
                # Predict time and cost using the appropriate model
                module_encoded = label_enc_module.transform([best_match.lower()])[0] if best_match.lower() in label_enc_module.classes_ else len(label_enc_module.classes_)
                lang_encoded = label_enc_lang.transform([technology.lower()])[0] if technology.lower() in label_enc_lang.classes_ else len(label_enc_lang.classes_)

                # Log encoding details for debugging
                logging.info(f"Module: '{cleaned_module_name}' -> Encoded: {module_encoded} (in classes: {cleaned_module_name.lower() in label_enc_module.classes_})")
                logging.info(f"Technology: '{technology}' -> Encoded: {lang_encoded} (in classes: {technology.lower() in label_enc_lang.classes_})")
                
                # Log available classes for debugging (only once)
                if module_encoded == len(label_enc_module.classes_):
                    logging.warning(f"Available module classes: {list(label_enc_module.classes_)[:10]}...")  # Show first 10
                if lang_encoded == len(label_enc_lang.classes_):
                    logging.warning(f"Available technology classes: {list(label_enc_lang.classes_)}")
                
                # Log the actual input to the model
                logging.info(f"Model input: module_encoded={module_encoded}, lang_encoded={lang_encoded}")
                
                # Predict time and cost
                predicted_time = round(user_time_model.predict([[module_encoded, lang_encoded]])[0], 2)
                predicted_cost = round(user_cost_model.predict([[module_encoded, lang_encoded]])[0], 2)

                # Log predicted values for debugging
                logging.info(f"Predicted Time for {cleaned_module_name}: {predicted_time}, Predicted Cost: {predicted_cost}")
                
                # Check if predictions are too similar (indicates model issue)
                if hasattr(user_time_model, 'predict'):
                    # Test with different inputs to see if model is working
                    test_inputs = [
                        [0, 0],  # First module, first technology
                        [1, 1],  # Second module, second technology  
                        [module_encoded, lang_encoded]  # Current input
                    ]
                    test_predictions = [user_time_model.predict([input_val])[0] for input_val in test_inputs]
                    logging.info(f"Test predictions for different inputs: {test_predictions}")
                    if len(set(test_predictions)) == 1:
                        logging.warning(" Model is predicting same values for different inputs - possible model issue!")

                complexity = "Small" if predicted_time <= 10 else "Medium" if predicted_time <= 30 else "Large"

                # Prepare the module data for saving
                module_data = {
                    "Module Name": cleaned_module_name,  # Use cleaned name
                    "Technology": technology,
                    "Time": predicted_time,
                    "Cost": predicted_cost,
                    "Complexity": complexity
                }

                # Check if module exists in user dataset and replace if it does
                module_found = False
                for existing_module in user_dataset:
                    if existing_module["Module Name"].lower() == module_name.lower() and existing_module["Technology"].lower() == technology.lower():
                        # Replace existing module data with new predictions
                        existing_module.update(module_data)
                        module_found = True
                        break

                # If module not found, add as a new entry
                if not module_found:
                    user_dataset.append(module_data)

            logging.info(f"Updated user dataset: {user_dataset}")

            # Save updated user dataset
            with open(user_file_path, "w") as f:
                json.dump(user_dataset, f, indent=4)

            # Train a new model if this is the second time the user is using the service
            if len(user_dataset) > 0:  # Assuming if there's data, it's the second time
                from model.new_train import main as train_model_main
                train_model_main(user_id)  # Call the function to train the model with user data
                logging.info(f"Trained new model for user {user_id}.")

        # Return the full response including client name and project title
        return jsonify({
            "modules": response["modules"],
            "client_name": response["client_name"],
            "project_title": response["project_title"]
        }), 200
    
    except Exception as e:
        logging.error(f"Error in module_details_api: {str(e)}")
        import traceback
        logging.error(f"Full traceback: {traceback.format_exc()}")
        
        if "429" in str(e):
            rotate_api_key()  # Rotate API key if rate limit error occurs
            return jsonify({"error": "Rate limit exceeded. Retrying with new API key."}), 429
        else:
            return jsonify({"error": f"An error occurred: {str(e)}"}), 500
    except Exception as e:
        logging.error(f"Error in module_details_api (outer): {str(e)}")
        import traceback
        logging.error(f"Full traceback: {traceback.format_exc()}")
        return jsonify({"error": f"An error occurred: {str(e)}"}), 500
    


# Create the blueprint for module routes
module_blueprint = Blueprint('module', __name__)

@module_blueprint.route("/update-module", methods=["POST"])
def update_module():
    """Update module details in dataset.json"""

    try:
        data = request.get_json()

        # Get user_id from top level
        user_id = data.get("user_id")
        if not user_id:
            return jsonify({"error": "Missing user_id at top level."}), 400

        # Expecting a list of module data
        modules_data = data.get("modules")

        if not modules_data or not isinstance(modules_data, list):
            return jsonify({"error": "Invalid data format. Expected a list of modules."}), 400

        for module_data in modules_data:
            # Map incoming fields to desired fields
            module_name = module_data.get("Module Name")  # Adjusted field name
            technology = module_data.get("Technology")  # Adjusted field name
            updated_time = module_data.get("Time")  # Adjusted field name
            updated_cost = module_data.get("Cost")  # Adjusted field name
            complexity = module_data.get("Complexity")  # Adjusted field name

            if not all([module_name, technology, updated_time, updated_cost]):
                return jsonify({"error": "Missing required fields in one of the modules."}), 400

            # Define user-specific directory and file path
            user_directory = os.path.join('user_data', str(user_id))
            os.makedirs(user_directory, exist_ok=True)
            user_file_path = os.path.join(user_directory, f'{user_id}.json')

            # Load existing user dataset
            try:
                with open(user_file_path, "r") as f:
                    user_dataset = json.load(f)
            except (FileNotFoundError, json.JSONDecodeError):
                user_dataset = []

            # Check if module exists in user dataset
            module_found = False
            for module in user_dataset:
                if module["Module Name"].lower() == module_name.lower() and module["Technology"].lower() == technology.lower():
                    module["Time"] = updated_time
                    module["Cost"] = updated_cost
                    module["Complexity"] = complexity
                    module_found = True
                    break

            # If module not found, add as a new entry
            if not module_found:
                user_dataset.append({
                    "Module Name": module_name,
                    "Technology": technology,
                    "Time": updated_time,
                    "Cost": updated_cost,
                    "Complexity": complexity
                })

            # Save updated user dataset
            with open(user_file_path, "w") as f:
                json.dump(user_dataset, f, indent=4)

        # After processing all modules, train the model for the specific user
        from model.new_train import main as train_main
        train_main(user_id=user_id)

        return jsonify({"message": "Module details updated and model retrained successfully!"}), 200

    except Exception as e:
        logging.error(f"Error updating module: {e}")
        return jsonify({"error": "An error occurred while updating module details"}), 500


